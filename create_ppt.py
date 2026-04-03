#!/usr/bin/env python3
"""Generate DeployDock Major Project Presentation (11 slides)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colour palette ──────────────────────────────────────────────
DARK_BG      = RGBColor(0x0F, 0x17, 0x2A)  # Deep navy
ACCENT_BLUE  = RGBColor(0x38, 0xBD, 0xF8)  # Bright cyan-blue
ACCENT_PURPLE= RGBColor(0xA7, 0x8B, 0xFA)  # Soft purple
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCB, 0xD5, 0xE1)
MID_GRAY     = RGBColor(0x94, 0xA3, 0xB8)
GREEN        = RGBColor(0x4A, 0xDE, 0x80)
ORANGE       = RGBColor(0xFB, 0xBF, 0x24)
RED_SOFT     = RGBColor(0xF8, 0x71, 0x71)
CARD_BG      = RGBColor(0x1E, 0x29, 0x3B)  # Slightly lighter navy

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helpers ─────────────────────────────────────────────────────
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width or Pt(1.5)
    else:
        shape.line.fill.background()
    # Smaller corner radius
    shape.adjustments[0] = 0.05
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_slide_content(slide, items, left, top, width, height,
                              font_size=16, color=LIGHT_GRAY, bullet_color=ACCENT_BLUE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(8)
        p.space_before = Pt(4)
        # Bullet character
        run_bullet = p.add_run()
        run_bullet.text = "\u25B8  "
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.name = "Segoe UI"
        # Text
        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = color
        run_text.font.name = "Segoe UI"
    return txBox

def add_section_title(slide, title, subtitle=None):
    """Standard section header with accent line."""
    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(0.6), Inches(0.15), Inches(0.6))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()
    line.rotation = 0

    add_text_box(slide, Inches(1.2), Inches(0.45), Inches(10), Inches(0.7),
                 title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(1.2), Inches(1.1), Inches(10), Inches(0.5),
                     subtitle, font_size=16, color=MID_GRAY)

def add_slide_number(slide, num, total=11):
    add_text_box(slide, Inches(12.0), Inches(7.0), Inches(1.2), Inches(0.4),
                 f"{num} / {total}", font_size=11, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)

def add_card(slide, left, top, width, height, title, items,
             accent=ACCENT_BLUE, icon_text=None):
    """A styled card with title and bullet items."""
    card = add_rect(slide, left, top, width, height, CARD_BG, border_color=accent, border_width=Pt(1.5))

    # Icon circle if provided
    if icon_text:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                         left + Inches(0.3), top + Inches(0.25),
                                         Inches(0.45), Inches(0.45))
        circle.fill.solid()
        circle.fill.fore_color.rgb = accent
        circle.line.fill.background()
        tf = circle.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = icon_text
        run.font.size = Pt(16)
        run.font.color.rgb = DARK_BG
        run.font.bold = True
        title_left = left + Inches(0.9)
    else:
        title_left = left + Inches(0.3)

    add_text_box(slide, title_left, top + Inches(0.2), width - Inches(1), Inches(0.45),
                 title, font_size=18, color=WHITE, bold=True)

    add_bullet_slide_content(slide, items,
                              left + Inches(0.3), top + Inches(0.75),
                              width - Inches(0.6), height - Inches(0.9),
                              font_size=13, color=LIGHT_GRAY, bullet_color=accent)

# ═══════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, DARK_BG)

# Top decorative gradient bar
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_BLUE
bar.line.fill.background()

# Project title
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.9),
             "DeployDock", font_size=54, color=ACCENT_BLUE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(1.85), Inches(11), Inches(0.6),
             "Self-Hosted Docker Container Management & Deployment Dashboard",
             font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Separator line
sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(4.5), Inches(2.7), Inches(4.3), Inches(0.03))
sep.fill.solid()
sep.fill.fore_color.rgb = ACCENT_PURPLE
sep.line.fill.background()

# Team heading
add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.4),
             "TEAM MEMBERS", font_size=14, color=MID_GRAY, bold=True,
             alignment=PP_ALIGN.CENTER)

team_members = ["Divija MV", "Aatmasree Srinivas", "Jagrathi KS"]
for i, name in enumerate(team_members):
    add_text_box(slide, Inches(1), Inches(3.5 + i * 0.48), Inches(11), Inches(0.45),
                 name, font_size=20, color=WHITE, alignment=PP_ALIGN.CENTER)

# Guide
add_text_box(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.4),
             "GUIDE / MENTOR", font_size=14, color=MID_GRAY, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.65), Inches(11), Inches(0.45),
             "Dr. Sumathi Pawar", font_size=20, color=ACCENT_PURPLE,
             bold=True, alignment=PP_ALIGN.CENTER)

# Date
add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.4),
             "28 March 2026", font_size=16, color=MID_GRAY,
             alignment=PP_ALIGN.CENTER)

# Bottom bar
bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.42), SLIDE_W, Inches(0.08))
bar2.fill.solid()
bar2.fill.fore_color.rgb = ACCENT_PURPLE
bar2.line.fill.background()


# ═══════════════════════════════════════════════════════════════
# SLIDE 2: INTRODUCTION
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_title(slide, "Introduction", "What is DeployDock?")
add_slide_number(slide, 2)

add_text_box(slide, Inches(1.2), Inches(1.8), Inches(11), Inches(0.8),
             "DeployDock is a self-hosted Docker PaaS dashboard that provides a clean web interface\n"
             "to manage containers, deploy applications, and monitor resources in real-time.",
             font_size=18, color=LIGHT_GRAY)

cards_data = [
    ("The Problem", [
        "Docker CLI is powerful but not beginner-friendly",
        "Paid platforms (Heroku, Railway) are expensive",
        "No free, visual tool for small teams / students",
    ], RED_SOFT, "!"),
    ("Our Solution", [
        "Web-based dashboard for full Docker management",
        "One-click deployments & app templates",
        "Real-time logs, metrics & network visualization",
    ], GREEN, "\u2713"),
    ("Target Users", [
        "Students learning containerization",
        "Developers on personal VPS / homelab",
        "Small teams needing free PaaS",
    ], ACCENT_PURPLE, "\u2605"),
]

for i, (title, items, accent, icon) in enumerate(cards_data):
    add_card(slide, Inches(0.8 + i * 4.1), Inches(3.1), Inches(3.8), Inches(3.6),
             title, items, accent=accent, icon_text=icon)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3: OBJECTIVES
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_title(slide, "Project Objectives")
add_slide_number(slide, 3)

objectives = [
    "Design and develop a web-based Docker container management dashboard",
    "Implement container lifecycle operations (create, start, stop, restart, delete) via REST APIs",
    "Provide real-time log streaming and resource metrics using Server-Sent Events (SSE)",
    "Build an interactive network topology visualization using D3.js force-directed graphs",
    "Enable automated CI/CD deployments via GitHub webhooks with language auto-detection",
    "Implement role-based access control (RBAC) with three user roles: Admin, Developer, Viewer",
    "Offer one-click deployment templates for 8 popular services (Ghost, Gitea, Redis, etc.)",
    "Make Docker internals (namespaces, cgroups, OverlayFS) visually educational in the UI",
]

for i, obj in enumerate(objectives):
    y = Inches(1.7 + i * 0.65)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(1.0), y, Inches(0.4), Inches(0.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(i + 1)
    run.font.size = Pt(14)
    run.font.color.rgb = DARK_BG
    run.font.bold = True

    add_text_box(slide, Inches(1.6), y - Inches(0.02), Inches(10.5), Inches(0.45),
                 obj, font_size=16, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4: SYSTEM ARCHITECTURE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_title(slide, "System Architecture", "Full-stack Next.js 15 with Docker integration")
add_slide_number(slide, 4)

# Architecture layers - horizontal layout
layers = [
    ("Frontend", ACCENT_BLUE, [
        "Next.js 15 App Router",
        "React 19 + TypeScript",
        "Tailwind CSS + Shadcn/ui",
        "Recharts + D3.js",
        "Framer Motion",
    ]),
    ("API Layer", ACCENT_PURPLE, [
        "Next.js Route Handlers",
        "REST API endpoints",
        "SSE streaming routes",
        "NextAuth.js v5 (JWT)",
        "Zod validation",
    ]),
    ("Backend Services", GREEN, [
        "dockerode SDK",
        "Prisma ORM + SQLite",
        "simple-git (clone/pull)",
        "bcryptjs (passwords)",
        "Buildpack detection",
    ]),
    ("Infrastructure", ORANGE, [
        "Docker Engine API",
        "/var/run/docker.sock",
        "Container runtime",
        "Bridge networks",
        "OverlayFS layers",
    ]),
]

for i, (title, accent, items) in enumerate(layers):
    x = Inches(0.6 + i * 3.15)
    # Header bar
    hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  x, Inches(2.0), Inches(2.9), Inches(0.55))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = accent
    hdr.line.fill.background()
    hdr.adjustments[0] = 0.15
    tf = hdr.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(16)
    run.font.color.rgb = DARK_BG
    run.font.bold = True

    # Card body
    add_card(slide, x, Inches(2.65), Inches(2.9), Inches(3.8),
             "", items, accent=accent)

    # Arrow between layers
    if i < 3:
        arrow_x = x + Inches(2.95)
        arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                      arrow_x, Inches(4.2), Inches(0.22), Inches(0.35))
        arr.fill.solid()
        arr.fill.fore_color.rgb = MID_GRAY
        arr.line.fill.background()


# ═══════════════════════════════════════════════════════════════
# SLIDE 5: KEY FEATURES (1)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_title(slide, "Key Features", "Container Management & Real-Time Monitoring")
add_slide_number(slide, 5)

features1 = [
    ("\U0001F4E6", "Container Management", ACCENT_BLUE, [
        "List all containers with live status badges",
        "Start, stop, restart, pause, unpause actions",
        "Container details modal (logs, metrics, env)",
        "Force remove containers (admin only)",
    ]),
    ("\U0001F4CA", "Real-Time Metrics", GREEN, [
        "CPU & memory usage line charts (Recharts)",
        "Network I/O monitoring (bytes in/out)",
        "SSE streams updating every 2 seconds",
        "Last 60 data points for performance",
    ]),
    ("\U0001F4DD", "Live Log Streaming", ACCENT_PURPLE, [
        "Color-coded stdout (green) / stderr (red)",
        "Auto-scroll with manual pause toggle",
        "Server-Sent Events for real-time delivery",
        "Terminal-style dark viewer UI",
    ]),
]

for i, (icon, title, accent, items) in enumerate(features1):
    add_card(slide, Inches(0.5 + i * 4.2), Inches(2.0), Inches(3.9), Inches(4.5),
             f"{icon}  {title}", items, accent=accent)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6: KEY FEATURES (2)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_title(slide, "Key Features", "Deployment, Networking & Templates")
add_slide_number(slide, 6)

features2 = [
    ("\U0001F310", "Network Topology", ACCENT_BLUE, [
        "D3.js force-directed interactive graph",
        "Containers as circles, networks as diamonds",
        "Color-coded by state (running/stopped/paused)",
        "Drag, zoom, hover for details",
    ]),
    ("\U0001F680", "GitHub Auto-Deploy", ORANGE, [
        "Webhook receiver with HMAC-SHA256 verification",
        "Auto-detect: Node.js, Python, PHP, Static",
        "Generate Dockerfile if missing (buildpacks)",
        "Hot-swap running containers on push",
    ]),
    ("\U0001F3AF", "One-Click Templates", ACCENT_PURPLE, [
        "8 pre-configured services ready to deploy",
        "Ghost, Gitea, Nextcloud, Uptime Kuma",
        "Vaultwarden, n8n, Redis, PostgreSQL",
        "Pre-set ports and environment variables",
    ]),
]

for i, (icon, title, accent, items) in enumerate(features2):
    add_card(slide, Inches(0.5 + i * 4.2), Inches(2.0), Inches(3.9), Inches(4.5),
             f"{icon}  {title}", items, accent=accent)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7: DATABASE & SECURITY
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_title(slide, "Database Design & Security", "Prisma ORM + SQLite | NextAuth.js RBAC")
add_slide_number(slide, 7)

# Left: Database schema
add_card(slide, Inches(0.5), Inches(2.0), Inches(5.8), Inches(5.0),
         "\U0001F5C4  Database Schema (Prisma + SQLite)", [
             "User  \u2014  id, email, password (bcrypt), role, name",
             "Deployment  \u2014  appName, imageTag, status, logs, triggeredBy",
             "AppConfig  \u2014  appName, envVars (JSON), domain, port, repoUrl",
             "Account  \u2014  NextAuth OAuth provider links",
             "Session  \u2014  JWT session tokens & expiry",
             "VerificationToken  \u2014  Email verification tokens",
         ], accent=ACCENT_BLUE)

# Right: RBAC table
rbac_card = add_rect(slide, Inches(6.7), Inches(2.0), Inches(6.0), Inches(5.0),
                      CARD_BG, border_color=ACCENT_PURPLE, border_width=Pt(1.5))

add_text_box(slide, Inches(7.0), Inches(2.15), Inches(5.5), Inches(0.45),
             "\U0001F512  Role-Based Access Control", font_size=18, color=WHITE, bold=True)

# Table header
roles_header = ["Permission", "Viewer", "Developer", "Admin"]
col_x = [Inches(7.0), Inches(9.3), Inches(10.5), Inches(11.7)]
for j, hdr_text in enumerate(roles_header):
    c = ACCENT_PURPLE if j > 0 else WHITE
    add_text_box(slide, col_x[j], Inches(2.75), Inches(1.2), Inches(0.35),
                 hdr_text, font_size=13, color=c, bold=True)

# Separator
sep_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(7.0), Inches(3.1), Inches(5.4), Inches(0.02))
sep_line.fill.solid()
sep_line.fill.fore_color.rgb = MID_GRAY
sep_line.line.fill.background()

rows = [
    ("View containers/logs", "\u2713", "\u2713", "\u2713"),
    ("Pull images", "\u2717", "\u2713", "\u2713"),
    ("Create containers", "\u2717", "\u2713", "\u2713"),
    ("Start/stop/restart", "\u2717", "\u2713", "\u2713"),
    ("Edit env variables", "\u2717", "\u2713", "\u2713"),
    ("Delete containers", "\u2717", "\u2717", "\u2713"),
    ("Manage users", "\u2717", "\u2717", "\u2713"),
]

for i, (perm, v, d, a) in enumerate(rows):
    y = Inches(3.25 + i * 0.48)
    add_text_box(slide, col_x[0], y, Inches(2.2), Inches(0.35),
                 perm, font_size=12, color=LIGHT_GRAY)
    for j, val in enumerate([v, d, a]):
        clr = GREEN if val == "\u2713" else RED_SOFT
        add_text_box(slide, col_x[j + 1] + Inches(0.2), y, Inches(0.5), Inches(0.35),
                     val, font_size=14, color=clr, alignment=PP_ALIGN.CENTER, bold=True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8: TECH STACK
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_title(slide, "Technology Stack", "Modern full-stack web technologies")
add_slide_number(slide, 8)

tech_categories = [
    ("Frontend", ACCENT_BLUE, [
        ("Next.js 15", "App Router framework"),
        ("React 19", "UI library"),
        ("TypeScript 5", "Type safety"),
        ("Tailwind CSS", "Utility-first styling"),
        ("Shadcn/ui", "50+ UI components"),
    ]),
    ("Visualization", ACCENT_PURPLE, [
        ("Recharts", "Real-time charts"),
        ("D3.js v7", "Network graph"),
        ("Framer Motion", "Animations"),
        ("Lucide React", "Icon system"),
        ("next-themes", "Dark mode"),
    ]),
    ("Backend", GREEN, [
        ("dockerode v4", "Docker SDK"),
        ("Prisma ORM", "Database layer"),
        ("NextAuth.js v5", "Authentication"),
        ("simple-git", "Git operations"),
        ("bcryptjs", "Password hashing"),
    ]),
    ("Infrastructure", ORANGE, [
        ("SQLite", "Embedded database"),
        ("Docker Engine", "Container runtime"),
        ("SSE", "Real-time streaming"),
        ("Zod", "Schema validation"),
        ("pnpm", "Package manager"),
    ]),
]

for i, (cat, accent, techs) in enumerate(tech_categories):
    x = Inches(0.4 + i * 3.2)
    # Category header
    hdr = add_rect(slide, x, Inches(2.0), Inches(3.0), Inches(0.5), accent)
    tf = hdr.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = cat
    run.font.size = Pt(15)
    run.font.color.rgb = DARK_BG
    run.font.bold = True

    for j, (name, desc) in enumerate(techs):
        y = Inches(2.7 + j * 0.85)
        pill = add_rect(slide, x + Inches(0.1), y, Inches(2.8), Inches(0.7), CARD_BG,
                         border_color=accent, border_width=Pt(1))
        add_text_box(slide, x + Inches(0.3), y + Inches(0.05), Inches(2.5), Inches(0.3),
                     name, font_size=14, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(0.3), y + Inches(0.35), Inches(2.5), Inches(0.3),
                     desc, font_size=11, color=MID_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9: EDUCATIONAL FEATURES (OS CONCEPTS)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_title(slide, "Educational Design", "Making Docker internals visible in the UI")
add_slide_number(slide, 9)

add_text_box(slide, Inches(1.2), Inches(1.8), Inches(11), Inches(0.5),
             "When a container boots, the UI reveals the 5 underlying OS concepts step-by-step:",
             font_size=16, color=LIGHT_GRAY)

os_concepts = [
    ("\U0001D7CF", "Linux Namespaces", "Isolates PID tree, network stack, mount points, hostname",
     ACCENT_PURPLE, "Provides process isolation \u2014 each container sees only its own processes"),
    ("\U0001D7D0", "Control Groups", "Limits CPU, memory, and I/O consumption per container",
     ORANGE, "Prevents any single container from monopolizing host resources"),
    ("\U0001D7D1", "OverlayFS", "Copy-on-write writable layer on top of read-only image layers",
     ACCENT_BLUE, "Efficient storage \u2014 shared base layers, only diffs stored per container"),
    ("\U0001D7D2", "Bridge Network", "Virtual ethernet interface with iptables NAT rules",
     GREEN, "Enables port binding and inter-container communication"),
    ("\U0001D7D3", "PID 1 Process", "Runs entrypoint command, handles SIGTERM graceful shutdown",
     RED_SOFT, "Container's main process \u2014 when it exits, the container stops"),
]

for i, (num, title, desc, accent, detail) in enumerate(os_concepts):
    y = Inches(2.5 + i * 0.95)
    # Step card
    card = add_rect(slide, Inches(0.8), y, Inches(11.7), Inches(0.8), CARD_BG,
                     border_color=accent, border_width=Pt(1.5))

    # Number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(1.0), y + Inches(0.15), Inches(0.5), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = accent
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(i + 1)
    run.font.size = Pt(16)
    run.font.color.rgb = DARK_BG
    run.font.bold = True

    add_text_box(slide, Inches(1.7), y + Inches(0.05), Inches(2.5), Inches(0.35),
                 title, font_size=16, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.7), y + Inches(0.4), Inches(4.5), Inches(0.35),
                 desc, font_size=12, color=MID_GRAY)
    add_text_box(slide, Inches(6.5), y + Inches(0.2), Inches(5.5), Inches(0.4),
                 detail, font_size=12, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10: IMPLEMENTATION METHODOLOGY
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_title(slide, "Implementation Methodology", "Agile phased approach with incremental delivery")
add_slide_number(slide, 10)

phases = [
    ("Phase 0", "Project Setup", "Environment, dependencies, Prisma schema, DB seed", ACCENT_BLUE),
    ("Phase 1", "Docker API", "Container CRUD, image pull, REST endpoints", ACCENT_BLUE),
    ("Phase 2", "Real-Time Streams", "SSE log streaming, live metrics charts", GREEN),
    ("Phase 3", "Authentication", "NextAuth.js, login page, RBAC middleware", ACCENT_PURPLE),
    ("Phase 4", "Network Graph", "D3.js topology visualization, /network page", ORANGE),
    ("Phase 5", "Auto-Deploy CI/CD", "GitHub webhooks, buildpacks, hot-swap", RED_SOFT),
    ("Phase 6", "Templates & Env", "8 app templates, env vars editor", ACCENT_BLUE),
    ("Phase 7", "Polish & Testing", "Error boundaries, responsive UI, testing", GREEN),
]

# Timeline layout
for i, (phase, title, desc, accent) in enumerate(phases):
    col = i % 4
    row = i // 4
    x = Inches(0.5 + col * 3.15)
    y = Inches(2.0 + row * 2.8)

    # Phase card
    card = add_rect(slide, x, y, Inches(2.9), Inches(2.4), CARD_BG,
                     border_color=accent, border_width=Pt(1.5))

    # Phase badge
    badge = add_rect(slide, x + Inches(0.15), y + Inches(0.15), Inches(1.1), Inches(0.4), accent)
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = phase
    run.font.size = Pt(12)
    run.font.color.rgb = DARK_BG
    run.font.bold = True

    add_text_box(slide, x + Inches(0.2), y + Inches(0.65), Inches(2.5), Inches(0.4),
                 title, font_size=16, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.1), Inches(2.5), Inches(1.1),
                 desc, font_size=12, color=MID_GRAY)

    # Arrow between cards in same row
    if col < 3 and i < len(phases) - 1:
        arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                      x + Inches(2.95), y + Inches(1.0),
                                      Inches(0.18), Inches(0.3))
        arr.fill.solid()
        arr.fill.fore_color.rgb = MID_GRAY
        arr.line.fill.background()


# ═══════════════════════════════════════════════════════════════
# SLIDE 11: CONCLUSION & FUTURE SCOPE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_title(slide, "Conclusion & Future Scope")
add_slide_number(slide, 11)

# Conclusion card
add_card(slide, Inches(0.5), Inches(2.0), Inches(6.0), Inches(4.8),
         "\u2705  Conclusion", [
             "Successfully designed a full-stack Docker management platform",
             "Demonstrates OS concepts: namespaces, cgroups, OverlayFS, networking",
             "Real-time monitoring with SSE (Server-Sent Events)",
             "Automated CI/CD pipeline with GitHub webhook integration",
             "Role-based security with three-tier access control",
             "Educational UI making Docker internals visible and understandable",
             "Built with modern web stack: Next.js 15, React 19, TypeScript",
         ], accent=GREEN)

# Future Scope card
add_card(slide, Inches(6.8), Inches(2.0), Inches(6.0), Inches(4.8),
         "\U0001F52D  Future Scope", [
             "Kubernetes cluster management support",
             "Multi-node Docker Swarm orchestration",
             "Custom domain mapping with SSL (Let's Encrypt)",
             "CI/CD pipeline with GitLab & Bitbucket support",
             "Resource usage alerts & email notifications",
             "Mobile-responsive PWA for on-the-go management",
             "Plugin system for community extensions",
         ], accent=ACCENT_PURPLE)

# Thank you bar
thank_bar = add_rect(slide, Inches(2.5), Inches(6.95), Inches(8.3), Inches(0.45), CARD_BG,
                      border_color=ACCENT_BLUE, border_width=Pt(1.5))
add_text_box(slide, Inches(2.5), Inches(6.95), Inches(8.3), Inches(0.45),
             "Thank You!", font_size=20, color=ACCENT_BLUE, bold=True,
             alignment=PP_ALIGN.CENTER)

# ── Save ────────────────────────────────────────────────────────
output_path = "/home/hx0r/Downloads/dashed-droplets-dashboard/DeployDock_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
